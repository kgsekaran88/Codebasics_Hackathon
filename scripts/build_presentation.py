#!/usr/bin/env python3
"""
Build PowerPoint deck for AtliQ leadership + Codebasics RPC #21 submission.

Covers: problem statement, solution, tech stack, methodology, insights, editorial rec.

Usage:
  python scripts/build_presentation.py
  python scripts/build_presentation.py --output deck/TN_2026_AtliQ_Briefing.pptx
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from src.presentation_charts import export_all_slide_images
from src.story_metrics import StoryMetrics, load_story_metrics

# Brand colors
BG_RGB = RGBColor(15, 20, 25)
ACCENT_RGB = RGBColor(59, 130, 246)
TEXT_RGB = RGBColor(241, 245, 249)
MUTED_RGB = RGBColor(148, 163, 184)
GREEN_RGB = RGBColor(52, 211, 153)


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_RGB
    return slide


def _title(slide, text: str, top: float = 0.35, size: int = 30, width: float = 9.2):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(width), Inches(0.9))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = TEXT_RGB


def _subtitle(slide, text: str, top: float = 1.05, size: int = 14):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(8.9), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = MUTED_RGB


def _body(slide, lines: list[str], top: float = 1.55, left: float = 0.55, width: float = 8.9, size: int = 15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(3.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size if line and not line.startswith("•") else size - 1)
        if line.startswith("•"):
            p.font.color.rgb = TEXT_RGB
            p.level = 0
        elif line == "":
            p.font.size = Pt(6)
        else:
            p.font.color.rgb = MUTED_RGB if line.endswith(":") else TEXT_RGB
        p.space_after = Pt(6)


def _two_col_insight(
    slide,
    bullets: list[str],
    image: Path,
    chart_top: float = 1.45,
    chart_width: float = 5.35,
    text_width: float = 3.75,
):
    _body(slide, bullets, top=1.5, left=0.55, width=text_width, size=13)
    if image.exists():
        slide.shapes.add_picture(
            str(image),
            Inches(4.15),
            Inches(chart_top),
            width=Inches(chart_width),
        )


def _notes(slide, notes: str):
    slide.notes_slide.notes_text_frame.text = notes


def _top_flow(metrics: StoryMetrics) -> str:
    if not metrics.top_sankey_flows:
        return "—"
    a, b, n = metrics.top_sankey_flows[0]
    return f"{a} (2021) → {b} (2026): {n} seats"


def build_deck(metrics: StoryMetrics, images: dict[str, Path], output: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    flow = _top_flow(metrics)
    reg_tvk = ", ".join(
        f"{r} ({n})" for r, n in sorted(metrics.tvk_by_region.items(), key=lambda x: -x[1])[:3]
    )

    # --- 1. Title ---
    s = _blank_slide(prs)
    _title(s, "Tamil Nadu Assembly 2026", top=1.0, size=38)
    sub = s.shapes.add_textbox(Inches(0.55), Inches(1.85), Inches(8.9), Inches(1.2))
    p = sub.text_frame.paragraphs[0]
    p.text = "Neutral election briefing for AtliQ Media"
    p.font.size = Pt(22)
    p.font.color.rgb = ACCENT_RGB
    _body(
        s,
        [
            "Codebasics Resume Project Challenge #21",
            "ECI public data · 234 constituencies · 2021 vs 2026",
        ],
        top=2.75,
        size=14,
    )
    foot = s.shapes.add_textbox(Inches(0.55), Inches(4.85), Inches(8.5), Inches(0.4))
    foot.text_frame.paragraphs[0].text = metrics.headline
    foot.text_frame.paragraphs[0].font.size = Pt(10)
    foot.text_frame.paragraphs[0].font.color.rgb = MUTED_RGB
    _notes(s, f"Open with headline: {metrics.headline}")

    # --- 2. Problem statement ---
    s = _blank_slide(prs)
    _title(s, "The challenge", size=28)
    _subtitle(s, "Hackathon problem statement · RPC #21")
    _body(
        s,
        [
            "Client: AtliQ Media needs ideas for a one-hour TV special on the 2026 TN Assembly election.",
            "",
            "Constraints (non-negotiable):",
            "• Election Commission of India data only — no news, exit polls, or social media",
            "• Neutral, descriptive tone — no predictions, causal claims, or party commentary",
            "• Pick 3 of 6 research threads and connect them into one clear story",
            "",
            "Deliverables: 5–7 min video pitch · 8–10 slide deck · public GitHub · optional dashboard",
        ],
        top=1.45,
        size=14,
    )
    _notes(
        s,
        "Frame the RPC as a producer briefing, not a political argument. "
        "Mention six research questions exist; we chose margins, flips, and TVK vote share.",
    )

    # --- 3. Our solution ---
    s = _blank_slide(prs)
    _title(s, "Our solution", size=28)
    _subtitle(s, "Data → insights → editorial run-of-show → interactive briefing")
    _body(
        s,
        [
            "1. Reproducible pipeline",
            "• Python ETL on ECI CSVs → AC-level metrics, Sankey flows, regional tables",
            "",
            "2. Three lead narratives (research threads addressed)",
            "• Margin of victory — fragmentation and close races (Q6)",
            "• Seat flips — 2021→2026 party transitions (Q2)",
            "• Vote share & TVK — new entrant scale by region (Q3)",
            "",
            "3. Interactive briefing dashboard",
            "• Editorial brief, statewide KPIs, charts, constituency explorer",
            "",
            "4. This deck + video script — explicit 60-minute show recommendation",
        ],
        top=1.4,
        size=13,
    )
    _notes(s, "Emphasize selection skill: depth on three threads, not all six.")

    # --- 4. Tech stack ---
    s = _blank_slide(prs)
    _title(s, "Technology stack", size=28)
    _subtitle(s, "Built for reproducibility and demo-ready storytelling")
    _body(
        s,
        [
            "Data & analytics",
            "• Python 3 · pandas · Jupyter notebook (`notebooks/01_build_and_explore.ipynb`)",
            "• Shared metrics module (`src/metrics.py`, `docs/METRIC_DEFINITIONS.md`)",
            "",
            "Backend",
            "• FastAPI — REST API over `data/processed/` CSVs",
            "",
            "Frontend",
            "• React · TypeScript · Vite · Apache ECharts",
            "",
            "Deliverables automation",
            "• `scripts/build_processed_data.py` · `build_presentation.py` · Playwright E2E",
            "",
            "Deployment",
            "• Local: uvicorn :8000 · optional ngrok for sharing",
        ],
        top=1.35,
        size=13,
    )
    _notes(s, "Point reviewers to GitHub README for one-command reproduction.")

    # --- 5. Data & methodology ---
    s = _blank_slide(prs)
    _title(s, "Data & methodology", size=28)
    _subtitle(s, "How we turned ECI files into comparable 2021 vs 2026 metrics")
    _body(
        s,
        [
            "Sources",
            "• `tn_2021_results.csv`, `tn_2026_results.csv`, `constituency_master.csv` (234 ACs)",
            "",
            "Join key: assembly constituency number (`ac_number`) — not constituency name",
            "",
            "Definitions",
            "• Valid votes = sum of candidate votes excluding NOTA",
            "• Winner = highest votes among non-NOTA candidates per AC",
            "• Margin % = (winner − runner-up) / valid votes × 100",
            "• Flip = normalized 2021 winner party ≠ normalized 2026 winner party",
            "",
            "Party labels normalized before seat and flip counts (see README mapping table)",
        ],
        top=1.35,
        size=13,
    )
    _notes(s, "Stress descriptive-only analysis. 2026 turnout blank in starter file.")

    # --- 6. Research questions answered ---
    s = _blank_slide(prs)
    _title(s, "Research questions — coverage map", size=26)
    _subtitle(s, "Hackathon brief: pick 3 of 6 · we go deep on Q6/Q2/Q3 and cover Q1/Q4/Q5 as supporting views")
    _body(
        s,
        [
            f"• Q1 Geography → flip rate {metrics.top_flip_region} {metrics.top_flip_pct}% vs {metrics.lowest_flip_region} {metrics.lowest_flip_pct}%",
            f"• Q2 Flips (LEAD) → {metrics.flips_norm}/{metrics.total_acs} ACs · {metrics.flips_pct}% of state",
            f"• Q3 Vote share (LEAD) → TVK {metrics.tvk_vote_share_2026}% · DMK 2021 {metrics.dmk_vote_share_2021}% (statewide + by-region)",
            "• Q4 Reserved → GEN/SC/ST flip rates, margin distribution, party × category table (2021 vs 2026)",
            "• Q5 Turnout → state record 2026 ~85.1% · per-AC delta chart in dashboard when ECI Form-20 is loaded",
            f"• Q6 Margins (LEAD) → avg {metrics.avg_margin_2021}% → {metrics.avg_margin_2026}% · {metrics.winner_under_35_2026} winners <35% · {metrics.under_5_margin_2026} races <5pp",
            "",
            f"Plus deep electoral-science measures: ENP {metrics.enp_2021}→{metrics.enp_2026} · Pedersen {metrics.pedersen} · Gallagher LSq {metrics.gallagher_2021}→{metrics.gallagher_2026}",
        ],
        top=1.4,
        size=12,
    )
    _notes(s, "Show that we go beyond the six required questions with structural electoral-science indices.")

    # --- 7. Insight — margins ---
    s = _blank_slide(prs)
    _title(s, "Q6 — Margin of victory (LEAD)", size=24)
    _subtitle(s, "Closer races, fewer landslides")
    _two_col_insight(
        s,
        [
            "What changed",
            f"• Mean margin: {metrics.avg_margin_2021}% → {metrics.avg_margin_2026}%",
            "",
            "What it means for TV",
            f"• {metrics.under_5_margin_2026} seats under 5-point margin",
            f"• {metrics.winner_under_35_2026} plurality winners (<35% share)",
            "",
            "Neutral framing",
            "• Describe distribution shift — do not explain voter motivation",
        ],
        images["margins"],
    )
    _notes(s, "Point to 2026 beeswarm sitting below 2021.")

    # --- 8. Insight — flips ---
    s = _blank_slide(prs)
    _title(s, "Q2 — Seat churn at scale (LEAD)", size=24)
    _subtitle(s, "Where seats moved between 2021 and 2026")
    _two_col_insight(
        s,
        [
            "Scale",
            f"• {metrics.flips_norm} of {metrics.total_acs} ACs changed winner",
            f"• Largest flow: {flow}",
            "",
            "Regional pattern",
            f"• Highest flip rate: {metrics.top_flip_region} ({metrics.top_flip_pct}%)",
            f"• Lowest: {metrics.lowest_flip_region} ({metrics.lowest_flip_pct}%)",
            "",
            "Neutral framing",
            "• Say constituencies changed hands — not voters switched allegiance",
        ],
        images["sankey"],
    )
    _notes(s, f"Top flows: {'; '.join(f'{a}→{b}({n})' for a,b,n in metrics.top_sankey_flows[:3])}")

    # --- 9. Insight — TVK statewide ---
    s = _blank_slide(prs)
    _title(s, "Q3 — Vote share statewide (LEAD)", size=24)
    _subtitle(s, "Where TVK landed at the state level — descriptive comparison")
    _two_col_insight(
        s,
        [
            "Statewide",
            f"• TVK: {metrics.tvk_seats} seats · ~{metrics.tvk_vote_share_2026}% valid vote share",
            f"• DMK 2021 baseline: {metrics.dmk_seats_2021} seats · ~{metrics.dmk_vote_share_2021}% share",
            "",
            "Regional seats won (TVK)",
            f"• {reg_tvk}",
            "",
            "Neutral framing",
            "• Report shares side-by-side — do not claim whose votes moved where",
        ],
        images["vote_share"],
    )
    _notes(s, "Descriptive comparison only.")

    # --- 10. Q3 by-region ---
    s = _blank_slide(prs)
    _title(s, "Q3 — Vote share by region", size=24)
    _subtitle(s, "Region-level composition shows where TVK is heavier or lighter")
    if "vote_share_region" in images and images["vote_share_region"].exists():
        slide_image_full = images["vote_share_region"]
        s.shapes.add_picture(str(slide_image_full), Inches(0.55), Inches(1.4), width=Inches(8.9))
    _notes(
        s,
        "Region-level vote share within each macro-region. Shows that TVK and the major Dravidian parties "
        "carry very different shares in Chennai Metro vs Delta. Pair on air with the regional seats chart.",
    )

    # --- 11. Q1 Geography ---
    s = _blank_slide(prs)
    _title(s, "Q1 — Geographic pattern", size=24)
    _subtitle(s, "Where flips and TVK seats cluster")
    _two_col_insight(
        s,
        [
            "234-tile mosaic",
            "• One square = one AC · equal visual weight",
            "",
            "Regional contrast",
            f"• Top flip rate: {metrics.top_flip_region} ({metrics.top_flip_pct}%)",
            f"• Lowest: {metrics.lowest_flip_region} ({metrics.lowest_flip_pct}%)",
            "",
            "On air",
            "• Pair with district map and 2021↔2026 toggle",
        ],
        images["mosaic"],
        chart_top=1.4,
    )
    _notes(s, "Good B-roll: zoom one region at a time.")

    # --- 12. Q4 Reserved ---
    s = _blank_slide(prs)
    _title(s, "Q4 — Reserved seats (GEN / SC / ST)", size=24)
    _subtitle(s, "Did reserved constituencies behave differently from general seats?")
    if "reserved" in images and images["reserved"].exists():
        s.shapes.add_picture(str(images["reserved"]), Inches(0.55), Inches(1.35), width=Inches(8.9))
    _notes(
        s,
        "Compare 2021 vs 2026 winner mix in each reserved category. SC and GEN flip rates are similar "
        "in the data; ST is just 2 seats so treat as illustrative.",
    )

    # --- 13. Q5 Turnout ---
    s = _blank_slide(prs)
    _title(s, "Q5 — Turnout", size=24)
    _subtitle(s, "Statewide record · per-AC delta requires ECI Form-20 ingestion")
    if "turnout" in images and images["turnout"].exists():
        s.shapes.add_picture(str(images["turnout"]), Inches(0.55), Inches(1.35), width=Inches(8.9))
    _notes(
        s,
        "State-record 85.1% statewide turnout in 2026 per ECI. Per-AC turnout is intentionally blank "
        "in the starter CSV; pipeline reads data/external/turnout_2026.csv once analyst loads Form-20 "
        "values, then the Top-20 turnout-increase chart populates automatically.",
    )

    # --- Deep insight 1 — electoral arithmetic ---
    s = _blank_slide(prs)
    _title(s, "Deep insight — electoral arithmetic", size=24)
    _subtitle(s, "Indices that summarise the structural shift between 2021 and 2026")
    _two_col_insight(
        s,
        [
            "Fragmentation",
            f"• Effective Number of Parties: {metrics.enp_2021} → {metrics.enp_2026}",
            f"• Multi-cornered races (top-2 < 70%): {metrics.multi_cornered_races} of {metrics.total_acs}",
            "",
            "Volatility & proportionality",
            f"• Pedersen index: {metrics.pedersen} — above 20 is 'high' in literature",
            f"• Gallagher LSq: {metrics.gallagher_2021} → {metrics.gallagher_2026}",
            "",
            "Neutral framing",
            "• These are descriptive electoral-science measures — no causal claim",
        ],
        images.get("enp", Path("/nonexistent")),
    )
    _notes(
        s,
        "ENP rose from 3.74 to 4.29 — the assembly contest moved from 'two-and-a-half party' to "
        "'three-and-a-half party' arithmetic on vote-share weights. Pedersen 36 indicates the system "
        "underwent the largest party-share churn between any two elections we measured. "
        "Lower Gallagher in 2026 means seats track votes more closely than 2021 — fewer 'wasted' votes.",
    )

    # --- Deep insight 2 — regional swing ---
    s = _blank_slide(prs)
    _title(s, "Deep insight — regional swing", size=24)
    _subtitle(s, "Where each party gained or lost vote share between cycles")
    if "swing" in images and images["swing"].exists():
        s.shapes.add_picture(str(images["swing"]), Inches(0.55), Inches(1.35), width=Inches(8.9))
    _notes(
        s,
        f"Highest party-share churn: {metrics.top_pedersen_region} (Pedersen {metrics.top_pedersen_value}). "
        f"Anti-incumbency varies sharply by region: Chennai Metro {metrics.chennai_anti_incumbency}% of seats "
        f"changed party, vs Delta {metrics.delta_anti_incumbency}%. "
        "Heatmap shows green = vote-share gain, red = loss in percentage points.",
    )

    # --- Deep insight 3 — representation gap ---
    s = _blank_slide(prs)
    _title(s, "Deep insight — representation gap", size=24)
    _subtitle(s, "First-past-the-post amplifies leading parties — by how much in 2026?")
    _two_col_insight(
        s,
        [
            "How it works",
            "• Vote share = % of valid votes statewide",
            "• Seat share = % of 234 ACs won",
            "• Gap = seat share − vote share (pp)",
            "",
            "What the data shows (2026)",
            f"• TVK gap: {metrics.tvk_representation_gap:+.1f} pp (leading-party bonus)",
            f"• BJP gap: {metrics.bjp_representation_gap:+.1f} pp (small parties penalised)",
            "",
            "On air",
            "• Use to explain why seat count alone overstates dominance",
        ],
        images.get("representation_gap", Path("/nonexistent")),
    )
    _notes(
        s,
        "Gallagher LSq is the headline disproportionality measure (lower = fairer). It fell from "
        "15.15 in 2021 to 8.74 in 2026 because TVK's vote share (~35%) maps more closely to its seat "
        "share (~46%) than DMK's 38% / 56% split in 2021. Both elections are first-past-the-post wins, "
        "but the gap is structurally smaller this cycle.",
    )

    # --- 14. Editorial recommendation ---
    s = _blank_slide(prs)
    _title(s, "Editorial recommendation — 60-minute show", size=24)
    _subtitle(s, "Required deliverable: what AtliQ should air, in what order")
    _body(
        s,
        [
            "Segment 1 (10 min) — Statewide arithmetic",
            "• KPIs, vote-share bars, mosaic · set up fragmentation headline",
            "",
            "Segment 2 (15 min) — Regional picture",
            "• District map, flip-rate bars · Chennai vs Delta contrast",
            "",
            "Segment 3 (10 min) — Churn & close races",
            "• Margin distribution, Sankey flows, closest-races table",
            "",
            "Segment 4 (5 min) — Sources & caveats",
            "• ECI scope, party normalization, missing 2026 turnout",
            "",
            "Producer tool: constituency explorer for ad-hoc seat lookups during the show",
        ],
        top=1.3,
        size=13,
    )
    _notes(s, "Scored heavily in rubric. State minutes and order clearly.")

    # --- 15. Limitations & close ---
    s = _blank_slide(prs)
    _title(s, "Limitations & reproducibility", size=26)
    _body(
        s,
        [
            "Limitations",
            "• 2026 turnout not in starter CSV — Q5 optional; 2021 turnout only in dashboard",
            "• Constituency-level data — no booth-level or voter-age tables",
            "• Party labels normalized — raw ECI strings differ slightly from seat counts",
            "• Descriptive only — no forecasts or causal explanations",
            "",
            "Reproduce this analysis",
            "• `python scripts/build_all_deliverables.py`",
            "• Dashboard: `uvicorn api.main:app --port 8000` (after `cd web && npm run build`)",
            "",
            "Thank you — questions on methodology and constituency drill-down welcome",
        ],
        top=1.35,
        size=14,
    )
    _notes(s, "Close with honesty on data gaps. Offer live dashboard demo.")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return output


def main():
    import os

    mpl_dir = ROOT / ".mplconfig"
    mpl_dir.mkdir(exist_ok=True)
    os.environ.setdefault("MPLCONFIGDIR", str(mpl_dir))

    parser = argparse.ArgumentParser(description="Build PowerPoint deck for TN 2026 RPC")
    parser.add_argument(
        "--output",
        type=Path,
        default=ROOT / "deck" / "TN_2026_AtliQ_Briefing.pptx",
    )
    parser.add_argument(
        "--assets",
        type=Path,
        default=ROOT / "outputs" / "presentation_assets",
    )
    args = parser.parse_args()

    metrics = load_story_metrics()
    images = export_all_slide_images(args.assets, metrics)
    out = build_deck(metrics, images, args.output)
    print(f"Saved deck: {out} (18 slides)")
    print(f"Slide images: {args.assets}")
    print(f"\nHeadline:\n{metrics.headline}")


if __name__ == "__main__":
    main()
